#!/usr/bin/env python3
"""
PPTX to Markdown Converter

Converts PowerPoint (.pptx) files to Markdown (.md) format.
Each slide becomes a level 1 header with its title, and content is structured
with level 2 headers based on the slide content structure.

Usage:
    python pptx_to_md.py <folder_path>
    python pptx_to_md.py <file1.pptx> [file2.pptx] [...]
"""

import argparse
import sys
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def get_slide_title(slide) -> str:
    """
    Extract the title from a slide.
    Tries to find a title placeholder first, then falls back to
    the first significant text element.
    """
    # Try to find title placeholder
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "placeholder_format"):
                pf = shape.placeholder_format
                if pf.type in (14, 15):  # 14=Header, 15=Title, 1=Title (Centered)
                    if hasattr(shape, "text") and shape.text.strip():
                        return shape.text.strip()

    # Fallback: find the first text element, often the title
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            # Usually titles are shorter and don't have newlines
            if len(text) < 100 and "\n" not in text:
                return text

    return "Untitled"


def extract_paragraphs(shape) -> List[Dict]:
    """
    Extract paragraphs from a text shape with formatting info.
    """
    paragraphs = []
    if not hasattr(shape, "text_frame"):
        return paragraphs

    text_frame = shape.text_frame
    for para_idx, paragraph in enumerate(text_frame.paragraphs):
        if not paragraph.text.strip():
            continue

        para_info = {
            "text": paragraph.text.strip(),
            "level": paragraph.level,
            "is_bullet": False,
        }

        # Check if it's a bullet by accessing the pPr element
        try:
            pPr = paragraph._element.pPr
            if pPr is not None:
                # Check for various bullet properties
                buChar = getattr(pPr, "buChar", None)
                buAutoNum = getattr(pPr, "buAutoNum", None)
                buBlip = getattr(pPr, "buBlip", None)
                if buChar is not None or buAutoNum is not None or buBlip is not None:
                    para_info["is_bullet"] = True
        except AttributeError:
            pass

        # Check font size for hierarchy
        runs = list(paragraph.runs)
        if runs and runs[0].font.size:
            para_info["font_size"] = runs[0].font.size.pt

        paragraphs.append(para_info)

    return paragraphs


def extract_table_data(shape) -> List[List[str]]:
    """
    Extract table data from a table shape.
    """
    if not hasattr(shape, "table"):
        return []

    table = shape.table
    data = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            row_data.append(cell.text.strip())
        data.append(row_data)
    return data


def sanitize_markdown(text: str) -> str:
    """
    Sanitize text for markdown - escape special characters.
    """
    # Basic sanitization - escape backslashes and pipes in tables
    text = text.replace("\\", "\\\\")
    return text


def paragraphs_to_markdown(paragraphs: List[Dict], base_level: int = 2) -> str:
    """
    Convert paragraphs to markdown with proper hierarchy.
    """
    if not paragraphs:
        return ""

    lines = []
    for para in paragraphs:
        text = sanitize_markdown(para["text"])
        if not text:
            continue

        # Determine if this looks like a header (no bullet, short, maybe larger font)
        is_header = not para["is_bullet"] and len(text) < 80

        if is_header:
            # Use as a sub-header
            header_level = base_level
            lines.append(f"{'#' * header_level} {text}")
        elif para["is_bullet"]:
            indent = "  " * para["level"]
            lines.append(f"{indent}* {text}")
        else:
            lines.append(text)

    return "\n".join(lines)


def table_to_markdown(table_data: List[List[str]]) -> str:
    """
    Convert table data to markdown table format.
    """
    if not table_data or not table_data[0]:
        return ""

    lines = []
    # Header row
    headers = [sanitize_markdown(cell).replace("|", "\\|") for cell in table_data[0]]
    lines.append("| " + " | ".join(headers) + " |")

    # Separator row
    lines.append("|" + "|".join(["---"] * len(headers)) + "|")

    # Data rows
    for row in table_data[1:]:
        cells = [sanitize_markdown(cell).replace("|", "\\|") for cell in row]
        lines.append("| " + " | ".join(cells) + " |")

    return "\n".join(lines)


def get_shape_description(shape) -> Optional[str]:
    """
    Get a description for non-text shapes like pictures, charts, etc.
    """
    # Build descriptions dict dynamically to handle missing enum values
    descriptions = {}

    # Try to add common shape types if they exist
    for attr_name in ["PICTURE", "CHART", "GROUP", "LINE", "OVAL",
                      "RECTANGLE", "ROUNDED_RECTANGLE", "AUTO_SHAPE"]:
        try:
            descriptions[getattr(MSO_SHAPE_TYPE, attr_name, None)] = {
                "PICTURE": "Image",
                "CHART": "Chart",
                "GROUP": "Grouped content",
                "LINE": "Line",
                "OVAL": "Oval/Circle",
                "RECTANGLE": "Rectangle/Box",
                "ROUNDED_RECTANGLE": "Rounded rectangle",
                "AUTO_SHAPE": "Shape",
            }.get(attr_name, "Shape")
        except AttributeError:
            pass

    desc = descriptions.get(shape.shape_type, f"Shape type {shape.shape_type}")

    if hasattr(shape, "name") and shape.name:
        desc += f" ({shape.name})"

    return desc


def process_slide(slide, slide_num: int) -> str:
    """
    Process a single slide and return its markdown content.
    """
    title = get_slide_title(slide)
    md_lines = [f"\n# Slide {slide_num} - {title}\n"]

    # Group shapes by their position/type for better organization
    title_shape_found = False
    content_sections = []
    tables = []
    other_shapes = []

    for shape in slide.shapes:
        # Skip shapes with no text or table (we'll handle them separately)
        has_content = (
            (hasattr(shape, "text") and shape.text.strip()) or
            (hasattr(shape, "table"))
        )

        if not has_content:
            desc = get_shape_description(shape)
            if desc:
                other_shapes.append(desc)
            continue

        # Check if this is the title shape (skip to avoid duplication)
        if not title_shape_found and hasattr(shape, "text"):
            shape_text = shape.text.strip()
            if shape_text == title or (len(shape_text) < 100 and shape_text in title):
                title_shape_found = True
                continue

        # Extract content based on shape type
        if hasattr(shape, "table"):
            table_data = extract_table_data(shape)
            if table_data:
                tables.append(table_data)
        else:
            paragraphs = extract_paragraphs(shape)
            if paragraphs:
                # Determine section title based on first paragraph
                section_title = None
                for para in paragraphs:
                    if not para["is_bullet"] and len(para["text"]) < 50:
                        section_title = para["text"]
                        para["used_as_header"] = True
                        break

                content_sections.append({
                    "title": section_title,
                    "paragraphs": paragraphs
                })

    # Add content sections
    for section in content_sections:
        if section["title"]:
            md_lines.append(f"\n## {sanitize_markdown(section['title'])}")

        # Filter out paragraphs used as headers
        content_paras = [p for p in section["paragraphs"]
                         if not p.get("used_as_header")]

        if content_paras:
            md_lines.append(paragraphs_to_markdown(content_paras, base_level=3))
        md_lines.append("")

    # Add tables
    for table_data in tables:
        md_lines.append("\n## Data Table\n")
        md_lines.append(table_to_markdown(table_data))
        md_lines.append("")

    # Add other shapes as notes
    if other_shapes:
        md_lines.append("\n## Visual Elements\n")
        for desc in set(other_shapes):
            md_lines.append(f"- {desc}")
        md_lines.append("")

    return "\n".join(md_lines)


def process_pptx_file(pptx_path: Path, output_dir: Path) -> bool:
    """
    Process a single PPTX file and convert it to markdown.

    Returns True if successful, False otherwise.
    """
    try:
        print(f"Processing: {pptx_path.name}")

        prs = Presentation(str(pptx_path))
        md_content = []

        # Add document title
        md_content.append(f"# {pptx_path.stem}\n")
        md_content.append(f"*Converted from {pptx_path.name}*\n")
        md_content.append("---\n")

        # Process each slide
        for idx in range(len(prs.slides)):
            slide = prs.slides[idx]
            slide_md = process_slide(slide, idx + 1)
            md_content.append(slide_md)

        # Write to markdown file
        md_filename = pptx_path.stem + ".md"
        md_path = output_dir / md_filename

        with open(md_path, "w", encoding="utf-8") as f:
            f.write("\n".join(md_content))

        print(f"  -> Created: {md_filename}")
        return True

    except Exception as e:
        print(f"  -> Error: {e}", file=sys.stderr)
        return False


def main():
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint files to Markdown",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s ./reports/                    # Process all PPTX in folder
  %(prog)s report1.pptx report2.pptx     # Process specific files
        """
    )

    parser.add_argument(
        "paths",
        nargs="+",
        help="Folder containing PPTX files, or specific PPTX files to process"
    )

    parser.add_argument(
        "-o", "--output",
        type=Path,
        default=None,
        help="Output directory (default: same as input files)"
    )

    args = parser.parse_args()

    # Collect all PPTX files
    pptx_files = []

    for path_str in args.paths:
        path = Path(path_str)

        if not path.exists():
            print(f"Warning: Path does not exist: {path}", file=sys.stderr)
            continue

        if path.is_file():
            if path.suffix.lower() == ".pptx":
                pptx_files.append(path)
            else:
                print(f"Warning: Not a PPTX file: {path}", file=sys.stderr)

        elif path.is_dir():
            pptx_files.extend(path.glob("*.pptx"))
            pptx_files.extend(path.glob("*.PPTX"))

    if not pptx_files:
        print("No PPTX files found to process.", file=sys.stderr)
        return 1

    # Determine output directory
    if args.output:
        output_dir = args.output
        output_dir.mkdir(parents=True, exist_ok=True)
    else:
        # Use the directory of the first file
        output_dir = pptx_files[0].parent

    print(f"Found {len(pptx_files)} PPTX file(s)")
    print(f"Output directory: {output_dir}\n")

    # Process files
    success_count = 0
    for pptx_file in pptx_files:
        if process_pptx_file(pptx_file, output_dir):
            success_count += 1

    print(f"\nCompleted: {success_count}/{len(pptx_files)} files converted")

    return 0 if success_count == len(pptx_files) else 1


if __name__ == "__main__":
    sys.exit(main())
